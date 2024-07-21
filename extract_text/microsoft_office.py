"""
Functions and CLI for converting Microsoft Office files 
into simple ".txt" files.

The goal is to maintain all context.
"""
import os
import io

from docx import Document
from PIL import Image
from pptx import Presentation

from utils import image_object_to_str, change_extension


def extract_text_from_shape(shape) -> str:
    """Extract text from a shape pulled from a slide"""
    if shape.has_text_frame:
        return shape.text_frame.text
    return ''


def extract_notes_from_slide(slide) -> str:
    """Extract notes from a slide"""
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        return notes_slide.notes_text_frame.text
    return ''


def docx_to_str(docx_path: str) -> str:
    """
    Convert a .docx file to a string of text

    Parameters:
    docx_path (str): Path to the input .docx file.
    """
    # Load the .docx file
    doc = Document(docx_path)

    text = ''

    # Iterate over paragraphs and embedded images in the .docx file
    for i, para in enumerate(doc.paragraphs):
        text += f"Paragraph {i + 1}:\n{para.text}\n\n"
    
    # Iterate over parts to parse text from images
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            img = rel.target_part.blob
            image = Image.open(io.BytesIO(img))
            text += f"Image OCR Text:\n{image_object_to_str(image)}\n\n"
    
    return text


def docx_to_txt(docx_path: str, txt_path: str = ''):
    """
    Convert a .docx file to a .txt file.

    Parameters:
    docx_path (str): Path to the input .docx file.
    txt_path (str): Path to the output .txt file.
    """
    if not txt_path:
        txt_path = change_extension(docx_path, 'txt')
    
    # Open the .txt file for writing
    with open(txt_path, 'w', encoding='utf-8') as txt_file:
        txt_file.write(docx_to_str(docx_path))

    print(f"Converted {docx_path} to {txt_path}")


def pptx_to_str(pptx_path: str) -> str:
    """
    Convert a .pptx file to a string of text.

    Parameters:
    pptx_path (str): Path to the input .pptx file.
    txt_path (str): Path to the output .txt file.
    """

    text = ''

    # Load the .pptx file
    presentation = Presentation(pptx_path)

    # Iterate over slides
    for slide_number, slide in enumerate(presentation.slides, start=1):
        text += f"Slide {slide_number}:\n"
        
        # Extract text from slide shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text += f"Text from slide {slide_number}:\n{paragraph.text}\n\n"

        # Extract notes
        notes = extract_notes_from_slide(slide)
        if notes:
            text += f"Notes from slide {slide_number}:\n{notes}\n\n"

        # Extract images and run OCR
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                img = shape.image.blob
                image_obj = Image.open(io.BytesIO(img))
                img_text = image_object_to_str(image_obj)
                text += f"OCR Text from image in slide {slide_number}:\n{img_text}\n\n"
        
        text += '\n' + '-'*40 + '\n\n'

    return text


def pptx_to_txt(pptx_path: str, txt_path: str = ''):
    """
    Convert a .pptx file to a .txt file.

    Parameters:
    pptx_path (str): Path to the input .pptx file.
    txt_path (str): Path to the output .txt file.
    """
    if not txt_path:
        txt_path = change_extension(pptx_path, 'txt')
    
    # Open the .txt file for writing
    with open(txt_path, 'w', encoding='utf-8') as txt_file:
        txt_file.write(pptx_to_str(pptx_path))

    print(f"Converted {pptx_path} to {txt_path}")


if __name__ == "__main__":

    import argparse
    parser = argparse.ArgumentParser(description="Convert .docx or .pptx file to .txt file")
    parser.add_argument("--input", type=str, required=True, help="Path to the .docx or .pptx file")
    parser.add_argument("--output", type=str, help="Path to the output .txt file")
    args = parser.parse_args()

    if args.input.lower().endswith('.docx'):
        docx_to_txt(args.input, args.output)
    elif args.input.lower().endswith('.pptx'):
        pptx_to_txt(args.input, args.output)
    else:
        print("Invalid file type. Please provide a .docx or .pptx file.")
    