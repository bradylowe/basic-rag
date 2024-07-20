import os

from docx import Document
from pptx import Presentation


def extract_text_from_shape(shape):
    """Extract text from a shape"""
    if shape.has_text_frame:
        return shape.text_frame.text
    return ''


def extract_notes_from_slide(slide):
    """Extract notes from a slide"""
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        return notes_slide.notes_text_frame.text
    return ''


def docx_to_txt(docx_path, txt_path):
    """
    Convert a .docx file to a .txt file.

    Parameters:
    docx_path (str): Path to the input .docx file.
    txt_path (str): Path to the output .txt file.
    """
    # Load the .docx file
    doc = Document(docx_path)

    # Open the .txt file for writing
    with open(txt_path, 'w', encoding='utf-8') as txt_file:
        # Iterate over paragraphs in the .docx file and write them to the .txt file
        for para in doc.paragraphs:
            txt_file.write(para.text + '\n')

    print(f"Converted {docx_path} to {txt_path}")


def pptx_to_txt(pptx_path, txt_path):
    """
    Convert a .pptx file to a .txt file.

    Parameters:
    pptx_path (str): Path to the input .pptx file.
    txt_path (str): Path to the output .txt file.
    """
    # Load the .pptx file
    presentation = Presentation(pptx_path)

    # Open the .txt file for writing
    with open(txt_path, 'w', encoding='utf-8') as txt_file:
        # Iterate over slides
        for slide_number, slide in enumerate(presentation.slides, start=1):
            txt_file.write(f"Slide {slide_number}:\n")
            
            # Extract text from slide shapes
            for shape in slide.shapes:
                text = extract_text_from_shape(shape)
                if text:
                    txt_file.write(text + '\n')

            # Extract notes
            notes = extract_notes_from_slide(slide)
            if notes:
                txt_file.write(f"Notes:\n{notes}\n")
            
            txt_file.write('\n' + '-'*40 + '\n\n')

    print(f"Converted {pptx_path} to {txt_path}")


def calculate_output_txt_name(file_path: str, output_dir: str) -> str:
    """Calculate the output .txt file name based on the input file path."""
    name = os.path.splitext(file_path)[0]
    if output_dir:
        name = os.path.basename(file_path)
        return os.path.join(output_dir, f"{name}.txt")
    else:
        return f"{name}.txt"
        

def convert_file_to_txt(file_path: str, convert_type_to_txt, output_dir: str = ''):
    """Convert a single file to .txt format using the appropriate converter."""
    ext = file_path.split('.')[-1].lower()
    if ext in convert_type_to_txt:
        txt_path = calculate_output_txt_name(file_path, output_dir)
        if os.path.exists(txt_path):
            print(f"Skipping {txt_path}. File already exists.")
            return
        
        os.makedirs(os.path.dirname(txt_path), exist_ok=True)
        converter = convert_type_to_txt[ext]
        converter(file_path, txt_path)
        print(f"Converted {file_path} to {txt_path}")


def convert_files_to_txt(input_dir: str, convert_type_to_txt, output_dir: str= '', recursive: bool = False):
    """Convert all supported files in a directory to .txt format."""
    for root, _, files in os.walk(input_dir):
        for name in files:
            file_path = os.path.join(root, name)
            convert_file_to_txt(file_path, convert_type_to_txt, output_dir)
        if not recursive:
            break


if __name__ == "__main__":

    import argparse
    parser = argparse.ArgumentParser(description="Convert all supported files in a directory to .txt format")
    parser.add_argument("--input_dir", type=str, required=True, help="Path to the input directory")
    parser.add_argument("--output_dir", type=str, help="Path to the output directory")
    parser.add_argument("--recursive", action="store_true", help="Recursively explore the directory structure to convert all files")
    args = parser.parse_args()
    
    supported_types = {
        'docx': docx_to_txt,
        'pptx': pptx_to_txt,
    }

    convert_files_to_txt(args.input_dir, supported_types, args.output_dir, recursive=args.recursive)
